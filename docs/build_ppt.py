"""Genera la presentación ejecutiva en .pptx.

    python docs/build_ppt.py

Las 6 slides mapean 1:1 con lo que pide el correo de sustentación:
solución · enfoque · arquitectura · demostración · decisiones · valor y límites.

Se genera por código y no a mano para que rehacerla tras un cambio sea gratis, y
para que las cifras salgan de un solo sitio (`DATOS`) en vez de estar repetidas
en varias cajas de texto.

Requiere `python-pptx` (solo para esto; no entra en requirements.txt).
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SALIDA = Path(__file__).resolve().parent / "Content-Suite - Sustentacion.pptx"

# ── Paleta: neutros con sesgo frío hacia el acento; teal de instrumentación ──
INK = RGBColor(0x10, 0x16, 0x19)
INK_SOFT = RGBColor(0x46, 0x56, 0x5C)
INK_FAINT = RGBColor(0x74, 0x86, 0x8C)
PAPER = RGBColor(0xFC, 0xFD, 0xFD)
SURFACE = RGBColor(0xF2, 0xF5, 0xF6)
HAIRLINE = RGBColor(0xD2, 0xDB, 0xDE)
ACCENT = RGBColor(0x0E, 0x6C, 0x74)
ACCENT_DK = RGBColor(0x0A, 0x4F, 0x55)
FAIL = RGBColor(0xA8, 0x27, 0x1F)
PASS = RGBColor(0x1A, 0x6F, 0x45)

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)  # 16:9

DATOS = {
    "app": "content-suite-web.onrender.com",
    "api": "content-suite-api.onrender.com/docs",
    "langfuse": "cloud.langfuse.com/project/cmszm13ah0az1ad0cp6p4suff",
    "tests": "63",
    "gen_s": "~80 s",
    "chunks": "~25",
    "dims": "1536",
}


# ─────────────────────────────────────────────────────────── helpers


def caja(slide, x, y, w, h, texto="", *, size=14, bold=False, color=INK,
         font=SANS, align=PP_ALIGN.LEFT, espacio=Pt(4), anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = espacio
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tf


def parrafo(tf, texto, *, size=13, bold=False, color=INK_SOFT, font=SANS,
            espacio=Pt(6), sangria=0):
    p = tf.add_paragraph()
    p.space_after = espacio
    p.level = sangria
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def rect(slide, x, y, w, h, relleno=SURFACE, borde=HAIRLINE, grosor=Pt(0.75)):
    from pptx.enum.shapes import MSO_SHAPE

    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = relleno
    if borde is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = borde
        s.line.width = grosor
    s.shadow.inherit = False
    return s


def nueva(prs, numero: str, eyebrow: str, titulo: str):
    """Diapositiva en blanco con el riel de medición y el encabezado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fondo = slide.background.fill
    fondo.solid()
    fondo.fore_color.rgb = PAPER

    # Riel izquierdo: las marcas son el motivo del proyecto (medir), no adorno.
    rect(slide, 0, 0, Inches(0.62), H, relleno=SURFACE, borde=None)
    caja(slide, Inches(0.14), Inches(0.42), Inches(0.5), Inches(0.3),
         numero, size=11, bold=True, color=ACCENT, font=MONO, align=PP_ALIGN.CENTER)
    for i in range(26):
        rect(slide, Inches(0.14), Inches(1.0) + Emu(int(Inches(0.24) * i)),
             Inches(0.34), Pt(0.75), relleno=HAIRLINE, borde=None)

    x = Inches(1.05)
    caja(slide, x, Inches(0.5), Inches(11), Inches(0.28),
         eyebrow.upper(), size=10.5, bold=True, color=ACCENT, font=MONO)
    caja(slide, x, Inches(0.85), Inches(11.6), Inches(0.9),
         titulo, size=30, bold=True, color=INK)
    return slide, x


def tarjeta(slide, x, y, w, h, titulo, cuerpo, *, destacada=False):
    rect(slide, x, y, w, h,
         relleno=SURFACE, borde=ACCENT if destacada else HAIRLINE,
         grosor=Pt(1.25) if destacada else Pt(0.75))
    tf = caja(slide, x + Inches(0.18), y + Inches(0.14), w - Inches(0.36), h - Inches(0.28),
              titulo, size=13, bold=True, color=INK)
    for linea in cuerpo:
        parrafo(tf, linea, size=11, color=INK_SOFT, espacio=Pt(4))
    return tf


def dato(slide, x, y, w, numero, etiqueta):
    rect(slide, x, y, w, Inches(1.0))
    caja(slide, x + Inches(0.16), y + Inches(0.14), w - Inches(0.3), Inches(0.42),
         numero, size=24, bold=True, color=ACCENT, font=MONO)
    caja(slide, x + Inches(0.16), y + Inches(0.62), w - Inches(0.3), Inches(0.3),
         etiqueta, size=10.5, color=INK_SOFT)


# ─────────────────────────────────────────────────────────── slides


def slide_1(prs):
    slide, x = nueva(prs, "01", "Reto técnico · IA Engineer", "Content Suite")

    caja(slide, x, Inches(1.75), Inches(10.8), Inches(1.0),
         "Cuando una compañía de consumo masivo lanza cientos de productos al año, el cuello "
         "de botella no es generar contenido: es que todo suene y se vea como la misma marca. "
         "Content Suite convierte las reglas de marca en un artefacto que las máquinas pueden "
         "leer, aplicar y auditar.",
         size=15, color=INK_SOFT)

    y = Inches(3.15)
    ancho = Inches(2.55)
    for i, (n, e) in enumerate([
        ("4/4", "módulos entregados"),
        ("3", "roles con vistas propias"),
        (DATOS["tests"], "tests, sin red"),
        (DATOS["gen_s"], "de brief a manual indexado"),
    ]):
        dato(slide, x + Emu(int(ancho + Inches(0.2)) * i), y, ancho, n, e)

    tf = caja(slide, x, Inches(4.55), Inches(11), Inches(0.3),
              "Entregables", size=12, bold=True, color=INK)
    parrafo(tf, f"App        {DATOS['app']}", size=11.5, color=INK_SOFT, font=MONO)
    parrafo(tf, f"API        {DATOS['api']}", size=11.5, color=INK_SOFT, font=MONO)
    parrafo(tf, f"Langfuse   {DATOS['langfuse']}", size=11.5, color=INK_SOFT, font=MONO)
    parrafo(tf, "Repo       github.com/PC-Nostades/Content-Suite", size=11.5,
            color=INK_SOFT, font=MONO)

    caja(slide, x, Inches(6.55), Inches(11), Inches(0.4),
         "FastAPI · React · Supabase + pgvector · OpenAI gpt-5.6-luna · LangGraph · Langfuse",
         size=11, color=INK_FAINT, font=MONO)


def slide_2(prs):
    slide, x = nueva(prs, "02", "Enfoque",
                     "Un manual en prosa no se puede auditar")

    caja(slide, x, Inches(1.8), Inches(11), Inches(0.8),
         "La tesis del proyecto: si las reglas de marca no son medibles, ningún sistema puede "
         "aplicarlas. Por eso cada regla se genera como un objeto con severity, modality y un "
         "check_hint — la instrucción de cómo verificarla.",
         size=14, color=INK_SOFT)

    y = Inches(2.95)
    rect(slide, x, y, Inches(5.3), Inches(2.5), relleno=SURFACE)
    tf = caja(slide, x + Inches(0.22), y + Inches(0.18), Inches(4.9), Inches(2.1),
              "Lo que un LLM escribe por defecto", size=12, bold=True, color=FAIL)
    for t in ['"el logo debe verse bien"', '"usar colores cálidos"',
              '"buena legibilidad"', '"fotos auténticas"']:
        parrafo(tf, t, size=12, color=INK_FAINT, espacio=Pt(9))

    rect(slide, x + Inches(5.6), y, Inches(5.9), Inches(2.5), relleno=SURFACE, borde=ACCENT)
    tf = caja(slide, x + Inches(5.82), y + Inches(0.18), Inches(5.5), Inches(2.1),
              "Lo que este sistema exige", size=12, bold=True, color=PASS)
    for t in ["el logo debe ocupar ≥ 8 % del ancho de la pieza",
              "#E8552D primario, máximo 40 % del área",
              "contraste texto/fondo ≥ 4.5:1 (WCAG AA)",
              "luz lateral natural, producto ≥ 25 % del encuadre"]:
        parrafo(tf, t, size=12, color=INK, espacio=Pt(9))

    y2 = Inches(5.72)
    rect(slide, x, y2, Inches(11.5), Inches(1.15), relleno=SURFACE, borde=ACCENT, grosor=Pt(1.5))
    tf = caja(slide, x + Inches(0.25), y2 + Inches(0.16), Inches(11), Inches(0.9),
              "Resultado real sobre una pieza con el logo dibujado al 4 %:",
              size=11, color=INK_FAINT)
    parrafo(tf, '"la pieza mide ~1080 px de ancho y el logo ~44 px; ocupa cerca del 4.1 %, '
                'por debajo del mínimo exigido de 8 %"', size=13, bold=True, color=INK)
    parrafo(tf, "citando visual.visual.el_logo_debe_ocupar_al_menos", size=10,
            color=INK_FAINT, font=MONO)


def slide_3(prs):
    slide, x = nueva(prs, "03", "Arquitectura",
                     "Monolito modular con una capa de IA compartida")

    y = Inches(1.8)
    w = Inches(3.7)
    gap = Inches(0.2)
    for i, (t, cuerpo) in enumerate([
        ("RENDER · STATIC SITE", [
            "React 19 · Vite · Tailwind v4",
            "Proxea /api/* desde el mismo origen:",
            "CORS desaparece. Y no duerme.",
        ]),
        ("RENDER · FASTAPI", [
            "modules/  auth · brand_dna",
            "          creative · governance",
            "ai/  llm · chunking · retrieval",
        ]),
        ("SUPABASE · POSTGRES 17", [
            "pgvector vector(1536) · HNSW",
            "RLS cerrando PostgREST",
            "7 tablas · migraciones idempotentes",
        ]),
    ]):
        px = x + Emu(int(w + gap) * i)
        rect(slide, px, y, w, Inches(1.6),
             relleno=SURFACE, borde=ACCENT if i == 1 else HAIRLINE,
             grosor=Pt(1.25) if i == 1 else Pt(0.75))
        tf = caja(slide, px + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(1.3),
                  t, size=10.5, bold=True, color=ACCENT_DK, font=MONO)
        for linea in cuerpo:
            parrafo(tf, linea, size=10.5, color=INK_SOFT, espacio=Pt(2))

    y = Inches(3.6)
    for i, (t, cuerpo) in enumerate([
        ("OPENAI", ["gpt-5.6-luna (texto y visión)", "text-embedding-3-small · 1536 dims",
                    "reasoning.effort = none"]),
        ("LANGGRAPH", ["Solo el Módulo II, donde hay", "un ciclo con decisión real",
                       "Los nodos llaman a nuestra capa llm"]),
        ("LANGFUSE", ["Qué contexto se recuperó del RAG", "Qué prompt se envió",
                      "Cuánto tardó cada auditoría"]),
    ]):
        px = x + Emu(int(w + gap) * i)
        rect(slide, px, y, w, Inches(1.6))
        tf = caja(slide, px + Inches(0.18), y + Inches(0.14), w - Inches(0.36), Inches(1.3),
                  t, size=10.5, bold=True, color=ACCENT_DK, font=MONO)
        for linea in cuerpo:
            parrafo(tf, linea, size=10.5, color=INK_SOFT, espacio=Pt(2))

    tf = caja(slide, x, Inches(5.5), Inches(11.5), Inches(1.5),
              "Por qué ai/ vive fuera de modules/", size=12, bold=True, color=INK)
    parrafo(tf, "El Módulo II necesita retrieval y el III generate_vision. Si esas piezas "
                "vivieran dentro de brand_dna/, los módulos siguientes importarían «hacia "
                "adentro» de otro módulo — el anti-patrón que obliga a refactorizar. Añadir un "
                "módulo es una carpeta, una línea en api.py y una fila en nav.ts.",
            size=12, color=INK_SOFT)


def slide_4(prs):
    slide, x = nueva(prs, "04", "Demostración de funcionamiento",
                     "El recorrido, y qué mirar en cada paso")

    pasos = [
        ("1 · Creador abre un manual",
         ["Espectro de voz, léxico con chips rojos/verdes,",
          "paleta con contraste WCAG calculado en cliente.",
          "Mirar: cada regla lleva su check_hint."]),
        ("2 · Creative Engine genera",
         ["El sistema recupera reglas ANTES de escribir.",
          "Mirar: los rule_id aplicados, y el bloque de",
          "violaciones que el guardrail corrigió."]),
        ("3 · Aprobador A revisa el texto",
         ["Misma app, distinta navegación y permisos.",
          "Mirar: el badge de rol y que el botón de crear",
          "marca ya no existe."]),
        ("4 · Aprobador B audita la imagen",
         ["Sube una pieza con el logo pequeño.",
          "Mirar: el dictamen mide 4.1 % y cita la regla",
          "exacta del manual que se incumple."]),
        ("5 · Langfuse",
         ["La traza de esa misma interacción: contexto",
          "recuperado, prompt enviado y latencias."]),
        ("Prueba de gobernanza en vivo",
         ["curl con el token del aprobador contra POST",
          "/brands devuelve 403. El RBAC no es la UI:",
          "se aplica en el servidor."]),
    ]

    w, hgt = Inches(3.7), Inches(1.62)
    for i, (t, cuerpo) in enumerate(pasos):
        col, fila = i % 3, i // 3
        px = x + Emu(int(w + Inches(0.2)) * col)
        py = Inches(1.85) + Emu(int(hgt + Inches(0.22)) * fila)
        tarjeta(slide, px, py, w, hgt, t, cuerpo, destacada=(i == 3))

    rect(slide, x, Inches(5.5), Inches(11.5), Inches(1.15),
         relleno=SURFACE, borde=ACCENT, grosor=Pt(1.25))
    tf = caja(slide, x + Inches(0.25), Inches(5.66), Inches(11), Inches(0.9),
              "Nota sobre la demo", size=11.5, bold=True, color=INK)
    parrafo(tf, "El backend corre en el free tier de Render y duerme tras 15 min: la primera "
                "petición puede tardar ~60 s. Por eso hay marcas ya sembradas — la demo del "
                "visor nunca depende de una llamada en vivo — y la sustentación se apoya en "
                "vídeo grabado.", size=11.5, color=INK_SOFT)


def slide_5(prs):
    slide, x = nueva(prs, "05", "Principales decisiones",
                     "Qué se decidió y por qué")

    decisiones = [
        ("Reglas como objetos con check_hint",
         ["Sin una instrucción de verificación medible, el",
          "Módulo III solo podría opinar sobre estética."]),
        ("Chunking por dominio, no por tamaño",
         ["Sin metadata limpia no hay pre-filtrado, y una",
          "consulta sobre el logo devolvería léxico."]),
        ("RAG para guía, código para lo duro",
         ["Un check de palabra prohibida necesita 100 % de",
          "recall. La búsqueda semántica no lo garantiza."]),
        ("LangGraph solo en el Módulo II",
         ["Ahí hay un ciclo con decisión. El fan-out del",
          "Módulo I lo expresa mejor asyncio.gather."]),
        ("Capa de proveedor intercambiable",
         ["Gemini limita a 20 peticiones AL DÍA por modelo.",
          "Migrar a OpenAI tocó un solo archivo."]),
        ("202 + polling en el Módulo I",
         ["80 s de generación no caben en un request que",
          "deba sobrevivir a un refresh del navegador."]),
    ]

    w, hgt = Inches(5.6), Inches(1.35)
    for i, (t, cuerpo) in enumerate(decisiones):
        col, fila = i % 2, i // 2
        px = x + Emu(int(w + Inches(0.28)) * col)
        py = Inches(1.85) + Emu(int(hgt + Inches(0.2)) * fila)
        tarjeta(slide, px, py, w, hgt, t, cuerpo, destacada=(i == 4))

    caja(slide, x, Inches(6.6), Inches(11.5), Inches(0.4),
         "Las diez decisiones completas, con su contexto, están documentadas en el README.",
         size=11, color=INK_FAINT)


def slide_6(prs):
    slide, x = nueva(prs, "06", "Valor y limitaciones",
                     "Qué resuelve hoy y qué falta para producción")

    tf = caja(slide, x, Inches(1.85), Inches(5.5), Inches(3.4),
              "Valor al negocio", size=14, bold=True, color=ACCENT)
    for t in [
        "El manual deja de ser un PDF que nadie abre: se vuelve reglas que el sistema aplica solo.",
        "La revisión de marca deja de depender de quién esté disponible. Mismo criterio en la pieza 1 y en la 200.",
        "Cada decisión es rastreable: un hallazgo cita su regla; una pieza cita las reglas que la guiaron.",
        "Riesgo regulatorio acotado: claims prohibidos y octógonos de la Ley 30021 son reglas del sistema.",
    ]:
        parrafo(tf, "·  " + t, size=12, color=INK_SOFT, espacio=Pt(10))

    tf = caja(slide, x + Inches(5.95), Inches(1.85), Inches(5.55), Inches(3.4),
              "Limitaciones, sin maquillar", size=14, bold=True, color=FAIL)
    for t in [
        "La auditoría visual estima proporciones, no mide píxeles. Acertó (44 px sobre 43 reales), pero cada hallazgo lleva confidence por eso.",
        "Cold start de ~60 s en el free tier. Mitigado con keep-alive y datos sembrados; no resuelto.",
        "Sin evaluación sistemática de calidad todavía: no hay golden set ni evals de regresión.",
        "El contenido generado requiere validación legal y nutricional. El sistema reduce el riesgo; no lo elimina.",
    ]:
        parrafo(tf, "·  " + t, size=12, color=INK_SOFT, espacio=Pt(10))

    y = Inches(5.5)
    rect(slide, x, y, Inches(11.5), Inches(1.35), relleno=SURFACE, borde=ACCENT, grosor=Pt(1.25))
    tf = caja(slide, x + Inches(0.25), y + Inches(0.16), Inches(11), Inches(1.1),
              "Lo que haría a continuación, en orden", size=12, bold=True, color=INK)
    parrafo(tf, "1 · Evals en Langfuse con manuales de referencia, para detectar regresión al "
                "tocar prompts.    2 · Medición geométrica real (OpenCV) para las reglas de logo "
                "y área, dejando al modelo de visión lo que sí es juicio.",
            size=11.5, color=INK_SOFT, espacio=Pt(3))
    parrafo(tf, "3 · Versionado de prompts en Langfuse para iterar sin redeploy.    "
                "4 · Cookie httpOnly + CSRF una vez validado el proxy de rewrite.",
            size=11.5, color=INK_SOFT)


def main() -> int:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    for fn in (slide_1, slide_2, slide_3, slide_4, slide_5, slide_6):
        fn(prs)

    prs.save(SALIDA)
    print(f"Presentación generada: {SALIDA}")
    print(f"  {len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas · 16:9")
    return 0


if __name__ == "__main__":
    sys.exit(main())
